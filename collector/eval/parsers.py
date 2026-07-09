"""
문서 타입별 텍스트 추출기 — pdf / docx / doc / xlsx / xls / csv / pptx / ppt.

각 파서는 실패 시 None을 반환한다 (샘플러가 스킵 처리).
"""
from __future__ import annotations

import csv
import re
import shutil
import signal
import subprocess
import tempfile
from pathlib import Path
from typing import Optional

_LIBREOFFICE = shutil.which("libreoffice") or shutil.which("soffice")


class _TimeoutError(Exception):
    pass


def _with_timeout(seconds: int):
    """일부 PDF/XLSX가 pdfplumber/openpyxl에서 무한정 걸리는 경우를 막기 위한
    SIGALRM 기반 타임아웃 데코레이터 (단일 스레드 전제, Unix 전용)."""
    def decorator(fn):
        def wrapper(*args, **kwargs):
            def _handler(signum, frame):
                raise _TimeoutError()
            old = signal.signal(signal.SIGALRM, _handler)
            signal.alarm(seconds)
            try:
                return fn(*args, **kwargs)
            except _TimeoutError:
                return None
            finally:
                signal.alarm(0)
                signal.signal(signal.SIGALRM, old)
        return wrapper
    return decorator


def _clean(text: str) -> str:
    text = re.sub(r"\n{3,}", "\n\n", text)
    text = re.sub(r"[ \t]{2,}", " ", text)
    return text.strip()


@_with_timeout(30)
def parse_pdf(path: Path, max_pages: int = 30) -> Optional[str]:
    try:
        import pdfplumber
    except ImportError:
        return None
    try:
        with pdfplumber.open(path) as pdf:
            parts = []
            for page in pdf.pages[:max_pages]:
                t = page.extract_text()
                if t and t.strip():
                    parts.append(t.strip())
        return _clean("\n\n".join(parts)) or None
    except Exception:
        return None


@_with_timeout(30)
def parse_docx(path: Path) -> Optional[str]:
    try:
        import docx
    except ImportError:
        return None
    try:
        d = docx.Document(str(path))
        parts = [p.text.strip() for p in d.paragraphs if p.text.strip()]
        for table in d.tables:
            for row in table.rows:
                cells = [c.text.strip() for c in row.cells if c.text.strip()]
                if cells:
                    parts.append("  ".join(cells))
        return _clean("\n".join(parts)) or None
    except Exception:
        return None


@_with_timeout(30)
def parse_pptx(path: Path) -> Optional[str]:
    try:
        from pptx import Presentation
    except ImportError:
        return None
    try:
        prs = Presentation(str(path))
        parts = []
        for i, slide in enumerate(prs.slides, 1):
            slide_lines = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    t = shape.text_frame.text.strip()
                    if t:
                        slide_lines.append(t)
                if shape.has_table:
                    for row in shape.table.rows:
                        cells = [c.text.strip() for c in row.cells if c.text.strip()]
                        if cells:
                            slide_lines.append("  ".join(cells))
            if slide_lines:
                parts.append(f"[슬라이드 {i}]\n" + "\n".join(slide_lines))
        return _clean("\n\n".join(parts)) or None
    except Exception:
        return None


@_with_timeout(30)
def parse_xlsx(path: Path, max_rows_per_sheet: int = 200) -> Optional[str]:
    try:
        import openpyxl
    except ImportError:
        return None
    try:
        wb = openpyxl.load_workbook(str(path), data_only=True, read_only=True)
        parts = []
        for ws in wb.worksheets:
            rows = []
            for i, row in enumerate(ws.iter_rows(values_only=True)):
                if i >= max_rows_per_sheet:
                    break
                vals = [str(v).strip() for v in row if v is not None and str(v).strip()]
                if vals:
                    rows.append("  ".join(vals))
            if rows:
                parts.append(f"=== {ws.title} ===\n" + "\n".join(rows))
        return _clean("\n\n".join(parts)) or None
    except Exception:
        return None


@_with_timeout(30)
def parse_xls(path: Path, max_rows_per_sheet: int = 200) -> Optional[str]:
    try:
        import xlrd
    except ImportError:
        return None
    try:
        wb = xlrd.open_workbook(str(path))
        parts = []
        for sh in wb.sheets():
            rows = []
            for r in range(min(sh.nrows, max_rows_per_sheet)):
                vals = [str(sh.cell(r, c).value).strip()
                        for c in range(sh.ncols) if str(sh.cell(r, c).value).strip()]
                if vals:
                    rows.append("  ".join(vals))
            if rows:
                parts.append(f"=== {sh.name} ===\n" + "\n".join(rows))
        return _clean("\n\n".join(parts)) or None
    except Exception:
        return None


def parse_csv(path: Path, max_rows: int = 300) -> Optional[str]:
    try:
        with path.open(encoding="utf-8-sig", errors="ignore", newline="") as f:
            rows = []
            for i, row in enumerate(csv.reader(f)):
                if i >= max_rows:
                    break
                vals = [v.strip() for v in row if v.strip()]
                if vals:
                    rows.append("  ".join(vals))
        return _clean("\n".join(rows)) or None
    except Exception:
        return None


def _libreoffice_to_text(path: Path) -> Optional[str]:
    """python-docx/python-pptx 등이 실패한 구형 포맷(.doc/.ppt/.odf/.rtf) 폴백."""
    if not _LIBREOFFICE:
        return None
    with tempfile.TemporaryDirectory() as tmp:
        try:
            subprocess.run(
                [_LIBREOFFICE, "--headless", "--convert-to", "txt:Text", "--outdir", tmp, str(path)],
                stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL, timeout=60, check=False,
            )
        except Exception:
            return None
        out = Path(tmp) / (path.stem + ".txt")
        if not out.exists():
            return None
        try:
            return _clean(out.read_text(encoding="utf-8", errors="ignore")) or None
        except Exception:
            return None


_PARSERS = {
    "pdf": parse_pdf,
    "docx": parse_docx,
    "pptx": parse_pptx,
    "xlsx": parse_xlsx,
    "xls": parse_xls,
    "csv": parse_csv,
}
# LibreOffice 변환 폴백이 필요한 구형/희귀 포맷
_LO_FALLBACK_EXTS = {"doc", "odf", "rtf", "ppt"}


def extract_text(path: Path) -> Optional[str]:
    """확장자 기반으로 적절한 파서를 골라 텍스트를 추출한다."""
    ext = path.suffix.lower().lstrip(".")
    if ext in _LO_FALLBACK_EXTS:
        return _libreoffice_to_text(path)
    parser = _PARSERS.get(ext)
    if parser is None:
        return None
    text = parser(path)
    if text:
        return text
    # 기본 파서가 실패하면(손상/구형 서브포맷) LibreOffice로 재시도
    return _libreoffice_to_text(path)
